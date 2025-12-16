"""
File Parsing Module.

This module handles the extraction of raw text from various file formats.
It abstracts away the specific libraries needed for PDFs, PowerPoint presentations,
and Images, providing a unified interface for the rest of the application.

Dependencies:
- PyPDF2: For PDF text extraction.
- python-pptx: For parsing slide decks.
- google-generativeai: For Optical Character Recognition (OCR) on images.
- PIL: For image file handling.
"""

import PyPDF2
from pptx import Presentation
from PIL import Image
import google.generativeai as genai

def parse_file(uploaded_file, api_key: str = None) -> str:
    """
    Parses a file object (PDF, PPTX, or Image) and returns its textual content.

    Args:
        uploaded_file: A Streamlit UploadedFile object.
        api_key (str, optional): Google Gemini API Key. Required ONLY for image parsing.

    Returns:
        str: The extracted text content, or an error message starting with "Error".
    """
    try:
        # Determine file type from extension
        # Note: We use the file extension as a heuristic for parsing strategy
        ext = uploaded_file.name.split('.')[-1].lower()
        text = ""

        # --- Strategy 1: PDF Parsing ---
        if ext == "pdf":
            reader = PyPDF2.PdfReader(uploaded_file)
            # Iterate through every page to concatenate text
            for page in reader.pages: 
                text += page.extract_text() or ""
                
        # --- Strategy 2: PowerPoint Parsing ---
        elif ext == "pptx":
            prs = Presentation(uploaded_file)
            # Iterate through slides -> shapes -> text frames
            for slide in prs.slides:
                for shape in slide.shapes:
                    # robustness check: not all shapes have text (e.g., lines, simple rects)
                    if hasattr(shape, "text"): 
                        text += shape.text + "\n"
                        
        # --- Strategy 3: Image OCR (via Gemini Vision) ---
        elif ext in ["jpg", "jpeg", "png"]:
            if not api_key: return "Error: API Key needed for vision."
            
            # Configure the Vision Model on the fly
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel("gemini-2.5-flash")
            
            # Load image into PIL format
            img = Image.open(uploaded_file)
            
            # Prompt the model to transcribe the image
            res = model.generate_content(["Extract text from this image:", img])
            text = res.text
            
        return text.strip()
        
    except Exception as e:
        # Return error as string to be handled by the UI
        return f"Error parsing file: {str(e)}"